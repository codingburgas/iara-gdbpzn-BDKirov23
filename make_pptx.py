from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

RED = RGBColor(0xDC, 0x35, 0x45)
DARK = RGBColor(0x0D, 0x11, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAD, 0xB5, 0xBD)
LIGHT_GRAY = RGBColor(0xE9, 0xEC, 0xEF)
GOLD = RGBColor(0xFF, 0xD7, 0x0A)
BLUE = RGBColor(0x1D, 0x35, 0x57)
GREEN = RGBColor(0x2A, 0x9D, 0x8F)
TEAL = RGBColor(0x45, 0x7B, 0x9D)
ORANGE = RGBColor(0xE7, 0x6F, 0x51)
LIME = RGBColor(0x52, 0xB7, 0x88)
CARD_BG = RGBColor(0x1A, 0x1A, 0x2E)

def add_bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tb

def add_bullet_text(slide, left, top, width, height, items, size=16, color=WHITE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "\u2022 " + item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return tb

def add_centered_box(slide, left, top, w, h, color, text, size=14):
    shape = add_rect(slide, left, top, w, h, color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return shape

# --- Slide 1: Title ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, Inches(0.3), H, RED)
add_text_box(slide, Inches(1), Inches(1.5), Inches(10), Inches(1.5),
    "\u0413\u0414\u041f\u0411\u0417\u041d \u2014 \u041e\u043f\u0435\u0440\u0430\u0442\u0438\u0432\u0435\u043d \u0446\u0435\u043d\u0442\u044a\u0440", size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3.2), Inches(10), Inches(1),
    "\u0421\u0438\u0441\u0442\u0435\u043c\u0430 \u0437\u0430 \u0443\u043f\u0440\u0430\u0432\u043b\u0435\u043d\u0438\u0435 \u043d\u0430 \u043f\u0440\u043e\u0438\u0437\u0448\u0435\u0441\u0442\u0432\u0438\u044f", size=28, color=GOLD)
add_text_box(slide, Inches(1), Inches(4.5), Inches(8), Inches(0.8),
    "\u0414\u043e\u043a\u043b\u0430\u0434\u0432\u0430\u043d\u0435 \u0438 \u043c\u043e\u043d\u0438\u0442\u043e\u0440\u0438\u043d\u0433  \u00b7  \u041a\u043e\u043e\u0440\u0434\u0438\u043d\u0430\u0446\u0438\u044f \u043d\u0430 \u0435\u043a\u0438\u043f\u0438  \u00b7  \u0410\u043d\u0430\u043b\u0438\u0437 \u0438 \u0441\u0442\u0430\u0442\u0438\u0441\u0442\u0438\u043a\u0430",
    size=16, color=GRAY)
add_rect(slide, Inches(1), Inches(5.8), Inches(3), Inches(0.04), RED)

# --- Slide 2: Problem ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, Inches(0.3), H, RED)
add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(1),
    "\u041f\u0440\u043e\u0431\u043b\u0435\u043c", size=40, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), RED)
add_bullet_text(slide, Inches(1), Inches(2), Inches(10), Inches(4.5), [
    "\u041b\u0438\u043f\u0441\u0430 \u043d\u0430 \u0446\u0435\u043d\u0442\u0440\u0430\u043b\u0438\u0437\u0438\u0440\u0430\u043d\u0430 \u0441\u0438\u0441\u0442\u0435\u043c\u0430 \u0437\u0430 \u0443\u043f\u0440\u0430\u0432\u043b\u0435\u043d\u0438\u0435 \u043d\u0430 \u043f\u0440\u043e\u0438\u0437\u0448\u0435\u0441\u0442\u0432\u0438\u044f",
    "\u0422\u0440\u0443\u0434\u043d\u0430 \u043a\u043e\u043e\u0440\u0434\u0438\u043d\u0430\u0446\u0438\u044f \u043c\u0435\u0436\u0434\u0443 \u0435\u043a\u0438\u043f\u0438 \u043d\u0430 \u0442\u0435\u0440\u0435\u043d \u0438 \u043e\u043f\u0435\u0440\u0430\u0442\u0438\u0432\u043d\u0438\u044f \u0446\u0435\u043d\u0442\u044a\u0440",
    "\u0411\u0430\u0432\u0435\u043d \u043e\u0431\u043c\u0435\u043d \u043d\u0430 \u0438\u043d\u0444\u043e\u0440\u043c\u0430\u0446\u0438\u044f \u043f\u0440\u0438 \u043a\u0440\u0438\u0442\u0438\u0447\u043d\u0438 \u0441\u0438\u0442\u0443\u0430\u0446\u0438\u0438",
    "\u041b\u0438\u043f\u0441\u0430 \u043d\u0430 \u0440\u0435\u0430\u043b\u043d\u043e\u0432\u0440\u0435\u043c\u0435\u043d\u043d\u0438 \u0434\u0430\u043d\u043d\u0438 \u0438 \u0441\u0442\u0430\u0442\u0438\u0441\u0442\u0438\u043a\u0430",
    "\u0420\u0430\u0437\u0447\u0438\u0442\u0430\u043d\u0435 \u043d\u0430 \u0442\u0435\u043b\u0435\u0444\u043e\u043d\u043d\u0438 \u043e\u0431\u0430\u0436\u0434\u0430\u043d\u0438\u044f \u0438 \u0445\u0430\u0440\u0442\u0438\u0435\u043d\u0438 \u0434\u043e\u043a\u043b\u0430\u0434\u0438",
], size=18, color=LIGHT_GRAY)

# --- Slide 3: Solution ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, Inches(0.3), H, RED)
add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(1),
    "\u041d\u0430\u0448\u0435\u0442\u043e \u0440\u0435\u0448\u0435\u043d\u0438\u0435", size=40, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), RED)
add_bullet_text(slide, Inches(1), Inches(2), Inches(10), Inches(4.5), [
    "\u0423\u0435\u0431-\u0431\u0430\u0437\u0438\u0440\u0430\u043d\u0430 \u043f\u043b\u0430\u0442\u0444\u043e\u0440\u043c\u0430 \u0434\u043e\u0441\u0442\u044a\u043f\u043d\u0430 \u043e\u0442 \u0432\u0441\u0435\u043a\u0438 \u0431\u0440\u0430\u0443\u0437\u044a\u0440",
    "REST API \u0438 Socket.IO \u0437\u0430 \u0440\u0435\u0430\u043b\u043d\u043e\u0432\u0440\u0435\u043c\u0435\u043d\u043d\u0430 \u043a\u043e\u043c\u0443\u043d\u0438\u043a\u0430\u0446\u0438\u044f",
    "JWT \u0443\u0434\u043e\u0441\u0442\u043e\u0432\u0435\u0440\u044f\u0432\u0430\u043d\u0435 \u0441 \u0440\u043e\u043b\u0438 (admin, dispatcher, commander, firefighter)",
    "\u0418\u043d\u0442\u0435\u0440\u0430\u043a\u0442\u0438\u0432\u043d\u0430 \u043a\u0430\u0440\u0442\u0430 \u0441 Google Maps / Leaflet",
    "AI \u0430\u0441\u0438\u0441\u0442\u0435\u043d\u0442 \u0441 Google Gemini \u0437\u0430 \u043f\u043e\u043c\u043e\u0449 \u043f\u0440\u0438 \u0432\u0437\u0435\u043c\u0430\u043d\u0435 \u043d\u0430 \u0440\u0435\u0448\u0435\u043d\u0438\u044f",
    "PWA \u2014 \u0438\u043d\u0441\u0442\u0430\u043b\u0438\u0440\u0430 \u0441\u0435 \u043d\u0430 \u0442\u0435\u043b\u0435\u0444\u043e\u043d \u043a\u0430\u0442\u043e native app",
], size=18, color=LIGHT_GRAY)

# --- Slide 4: Architecture ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, Inches(0.3), H, RED)
add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(1),
    "\u0410\u0440\u0445\u0438\u0442\u0435\u043a\u0442\u0443\u0440\u0430", size=40, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), RED)

arch_boxes = [
    (Inches(1), Inches(2), Inches(3.5), Inches(1.3), RED, "Flask REST API\nJWT \u043e\u0442\u043e\u0440\u0438\u0437\u0430\u0446\u0438\u044f"),
    (Inches(5), Inches(2), Inches(3.5), Inches(1.3), BLUE, "SQLite/PostgreSQL\nSQLAlchemy ORM"),
    (Inches(9), Inches(2), Inches(3.5), Inches(1.3), GREEN, "Socket.IO\nWebSocket"),
    (Inches(1), Inches(3.8), Inches(3.5), Inches(1.3), TEAL, "Google Gemini AI\n\u0427\u0430\u0442 \u0430\u0441\u0438\u0441\u0442\u0435\u043d\u0442"),
    (Inches(5), Inches(3.8), Inches(3.5), Inches(1.3), ORANGE, "Google Maps\n\u0418\u043d\u0442\u0435\u0440\u0430\u043a\u0442\u0438\u0432\u043d\u0430 \u043a\u0430\u0440\u0442\u0430"),
    (Inches(9), Inches(3.8), Inches(3.5), Inches(1.3), LIME, "PWA Service Worker\n\u041e\u0444\u043b\u0430\u0439\u043d \u043a\u0435\u0448"),
]
for left, top, w, h, color, text in arch_boxes:
    add_centered_box(slide, left, top, w, h, color, text)

# --- Slide 5: Features ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, Inches(0.3), H, RED)
add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(1),
    "\u041e\u0441\u043d\u043e\u0432\u043d\u0438 \u0444\u0443\u043d\u043a\u0446\u0438\u043e\u043d\u0430\u043b\u043d\u043e\u0441\u0442\u0438", size=40, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), RED)
add_bullet_text(slide, Inches(1), Inches(1.8), Inches(5.5), Inches(5), [
    "\u0422\u0430\u0431\u043b\u043e \u0441 \u043e\u043f\u0435\u0440\u0430\u0442\u0438\u0432\u043d\u0430 \u0438\u043d\u0444\u043e\u0440\u043c\u0430\u0446\u0438\u044f \u043d\u0430 \u0436\u0438\u0432\u043e",
    "\u0421\u044a\u0437\u0434\u0430\u0432\u0430\u043d\u0435 \u0438 \u043f\u0440\u043e\u0441\u043b\u0435\u0434\u044f\u0432\u0430\u043d\u0435 \u043d\u0430 \u043f\u0440\u043e\u0438\u0437\u0448\u0435\u0441\u0442\u0432\u0438\u044f",
    "\u0423\u043f\u0440\u0430\u0432\u043b\u0435\u043d\u0438\u0435 \u043d\u0430 \u0435\u043a\u0438\u043f\u0438 \u0438 \u0430\u0432\u0442\u043e\u043f\u0430\u0440\u043a",
    "\u0420\u0430\u0437\u043f\u0440\u0435\u0434\u0435\u043b\u0435\u043d\u0438\u0435 \u043d\u0430 \u0437\u0430\u0434\u0430\u0447\u0438 \u043a\u044a\u043c \u0435\u043a\u0438\u043f\u0438",
    "\u0421\u043c\u0435\u043d\u0435\u043d \u0433\u0440\u0430\u0444\u0438\u043a \u0438 \u043e\u0442\u043f\u0443\u0441\u043a\u0438",
], size=18, color=LIGHT_GRAY)
add_bullet_text(slide, Inches(6.5), Inches(1.8), Inches(5.5), Inches(5), [
    "\u0418\u043d\u0442\u0435\u0440\u0430\u043a\u0442\u0438\u0432\u043d\u0430 \u043a\u0430\u0440\u0442\u0430 \u0441 \u043f\u043e\u0437\u0438\u0446\u0438\u0438 \u0438 \u043c\u0430\u0440\u043a\u0435\u0440\u0438",
    "AI \u0447\u0430\u0442 \u0430\u0441\u0438\u0441\u0442\u0435\u043d\u0442 \u0441 Google Gemini",
    "\u0421\u0442\u0430\u0442\u0438\u0441\u0442\u0438\u043a\u0430 \u0441 Chart.js \u0433\u0440\u0430\u0444\u0438\u043a\u0438",
    "PDF \u043f\u0440\u043e\u0442\u043e\u043a\u043e\u043b\u0438 \u0437\u0430 \u043f\u0440\u043e\u0438\u0437\u0448\u0435\u0441\u0442\u0432\u0438\u044f",
    "SOS \u0431\u0443\u0442\u043e\u043d \u0437\u0430 \u0441\u043f\u0435\u0448\u043d\u0438 \u0441\u0438\u0433\u043d\u0430\u043b\u0438",
], size=18, color=LIGHT_GRAY)

# --- Slide 6: Dashboard ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, Inches(0.3), H, RED)
add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(1),
    "\u041e\u043f\u0435\u0440\u0430\u0442\u0438\u0432\u043d\u043e \u0442\u0430\u0431\u043b\u043e", size=40, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), RED)
add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(4),
    "\u0426\u0435\u043d\u0442\u0440\u0430\u043b\u0435\u043d \u0438\u0437\u0433\u043b\u0435\u0434 \u0441 4 \u043a\u043b\u044e\u0447\u043e\u0432\u0438 \u043f\u043e\u043a\u0430\u0437\u0430\u0442\u0435\u043b\u044f:\n"
    "\u2022 \u0410\u043a\u0442\u0438\u0432\u043d\u0438 \u043f\u0440\u043e\u0438\u0437\u0448\u0435\u0441\u0442\u0432\u0438\u044f\n"
    "\u2022 \u0415\u043a\u0438\u043f\u0438 \u043d\u0430 \u0442\u0435\u0440\u0435\u043d\n"
    "\u2022 \u041d\u0430\u043b\u0438\u0447\u043d\u0438 \u0441\u043b\u0443\u0436\u0438\u0442\u0435\u043b\u0438\n"
    "\u2022 \u041f\u0440\u0438\u043a\u043b\u044e\u0447\u0435\u043d\u0438 \u0434\u043d\u0435\u0441\n\n"
    "11 \u0441\u0435\u043a\u0446\u0438\u0438: \u0422\u0430\u0431\u043b\u043e, \u041f\u0440\u043e\u0438\u0437\u0448\u0435\u0441\u0442\u0432\u0438\u044f, \u0415\u043a\u0438\u043f\u0438,\n"
    "\u0421\u043c\u0435\u043d\u0438, \u041a\u0430\u0440\u0442\u0430, \u0417\u0430\u0434\u0430\u0447\u0438, \u0410\u0432\u0442\u043e\u043f\u0430\u0440\u043a, \u0421\u0442\u0430\u0442\u0438\u0441\u0442\u0438\u043a\u0430, AI \u0427\u0430\u0442, \u0420\u0435\u0441\u0443\u0440\u0441\u0438, API\n\n"
    "\u0422\u044a\u043c\u0435\u043d/\u0441\u0432\u0435\u0442\u044a\u043b \u0440\u0435\u0436\u0438\u043c. \u0410\u0434\u0430\u043f\u0442\u0438\u0432\u0435\u043d \u0434\u0438\u0437\u0430\u0439\u043d. Dark theme \u0438 PWA.",
    size=18, color=LIGHT_GRAY)

# --- Slide 7: Mobile PWA ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, Inches(0.3), H, RED)
add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(1),
    "\u041c\u043e\u0431\u0438\u043b\u043d\u043e \u043f\u0440\u0438\u043b\u043e\u0436\u0435\u043d\u0438\u0435 (PWA)", size=40, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), RED)
add_bullet_text(slide, Inches(1), Inches(2), Inches(10), Inches(4.5), [
    "\u0414\u043e\u0441\u0442\u044a\u043f\u043d\u043e \u043e\u0442 \u0432\u0441\u0435\u043a\u0438 \u0442\u0435\u043b\u0435\u0444\u043e\u043d \u2014 \u0431\u0435\u0437 \u0438\u043d\u0441\u0442\u0430\u043b\u0430\u0446\u0438\u044f \u043e\u0442 \u043c\u0430\u0433\u0430\u0437\u0438\u043d",
    "\u0414\u043e\u0431\u0430\u0432\u044f\u043d\u0435 \u043a\u044a\u043c \u043d\u0430\u0447\u0430\u043b\u043d\u0438\u044f \u0435\u043a\u0440\u0430\u043d (Android / iOS)",
    "Service Worker \u0437\u0430 \u043e\u0444\u043b\u0430\u0439\u043d \u0434\u043e\u0441\u0442\u044a\u043f",
    "Web App Manifest \u0441 \u0438\u043a\u043e\u043d\u0430 \u0438 \u0442\u0435\u043c\u0430",
    "\u0410\u0434\u0430\u043f\u0442\u0438\u0432\u0435\u043d \u0438\u043d\u0442\u0435\u0440\u0444\u0435\u0439\u0441 \u0437\u0430 \u043c\u0430\u043b\u043a\u0438 \u0435\u043a\u0440\u0430\u043d\u0438",
    "SOS \u0431\u0443\u0442\u043e\u043d \u0432\u0438\u043d\u0430\u0433\u0438 \u043d\u0430 \u0432\u0438\u0434\u043d\u043e \u043c\u044f\u0441\u0442\u043e",
], size=18, color=LIGHT_GRAY)

# --- Slide 8: API ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, Inches(0.3), H, RED)
add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(1),
    "REST API", size=40, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), RED)
add_bullet_text(slide, Inches(1), Inches(2), Inches(10), Inches(4.5), [
    "JWT Bearer token \u0437\u0430 \u0449\u0438\u0442\u0430 \u043d\u0430 endpoints",
    "18 SQLAlchemy \u043c\u043e\u0434\u0435\u043b\u0430 (User, Incident, Task, Team, FireTruck...)",
    "\u041f\u044a\u043b\u0435\u043d CRUD \u0437\u0430 \u043f\u0440\u043e\u0438\u0437\u0448\u0435\u0441\u0442\u0432\u0438\u044f, \u0435\u043a\u0438\u043f\u0438, \u0437\u0430\u0434\u0430\u0447\u0438, \u0440\u0435\u0441\u0443\u0440\u0441\u0438",
    "\u041a\u0430\u0447\u0432\u0430\u043d\u0435 \u043d\u0430 \u0441\u043d\u0438\u043c\u043a\u0438 \u043a\u044a\u043c \u043f\u0440\u043e\u0438\u0437\u0448\u0435\u0441\u0442\u0432\u0438\u044f \u0438 \u0430\u0432\u0442\u043e\u043c\u043e\u0431\u0438\u043b\u0438",
    "WebSocket \u0441\u0442\u0430\u0438 \u0437\u0430 real-time chat \u0438 emergency \u0441\u0438\u0433\u043d\u0430\u043b\u0438",
    "AI \u0447\u0430\u0442 endpoint \u0441 Google Gemini \u0438\u043d\u0442\u0435\u0433\u0440\u0430\u0446\u0438\u044f",
], size=18, color=LIGHT_GRAY)

# --- Slide 9: Tech Stack ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, Inches(0.3), H, RED)
add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(1),
    "\u0422\u0435\u0445\u043d\u043e\u043b\u043e\u0433\u0438\u0447\u0435\u043d \u0441\u0442\u0435\u043a", size=40, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), RED)

techs = [
    ("Backend", "Python 3, Flask, SQLAlchemy\nJWT, Socket.IO, Flask-CORS"),
    ("Database", "SQLite (dev) / PostgreSQL (prod)\nRedis (\u043e\u043f\u0446\u0438\u043e\u043d\u0430\u043b\u043d\u043e)"),
    ("Frontend", "HTML5, Vanilla JS, Bootstrap 5\nChart.js, jsPDF, Google Maps"),
    ("AI", "Google Gemini 2.0 Flash\nTogetherAI / HuggingFace (\u0440\u0435\u0437\u0435\u0440\u0432\u0430)"),
    ("Mobile", "PWA (Service Worker + Manifest)\n\u0418\u043d\u0441\u0442\u0430\u043b\u0430\u0446\u0438\u044f \u043e\u0442 \u0431\u0440\u0430\u0443\u0437\u044a\u0440"),
    ("DevOps", "Git, pip, Flask dev server\nGunicorn + Nginx (prod)"),
]
for i, (title, desc) in enumerate(techs):
    col = i % 3
    row = i // 3
    left = Inches(1 + col * 4)
    top = Inches(1.8 + row * 2.5)
    add_rect(slide, left, top, Inches(3.5), Inches(2), CARD_BG)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(3), Inches(0.6),
        title, size=20, color=GOLD, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.8), Inches(3), Inches(1),
        desc, size=14, color=LIGHT_GRAY)

# --- Slide 10: Demo ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, Inches(0.3), H, RED)
add_text_box(slide, Inches(1), Inches(2), Inches(10), Inches(1.5),
    "Live Demo", size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3.5), Inches(10), Inches(1),
    "http://localhost:5000", size=24, color=GOLD)
add_text_box(slide, Inches(1), Inches(4.5), Inches(10), Inches(1),
    "admin@gdpbzn.bg / admin123", size=18, color=GRAY)

# --- Slide 11: Thank you ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, Inches(0.3), H, RED)
add_text_box(slide, Inches(1), Inches(2), Inches(10), Inches(1),
    "\u0411\u043b\u0430\u0433\u043e\u0434\u0430\u0440\u044f \u0437\u0430 \u0432\u043d\u0438\u043c\u0430\u043d\u0438\u0435\u0442\u043e!", size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3.5), Inches(10), Inches(1),
    "\u0412\u044a\u043f\u0440\u043e\u0441\u0438?", size=28, color=GOLD)
add_rect(slide, Inches(1), Inches(4.8), Inches(3), Inches(0.04), RED)
add_text_box(slide, Inches(1), Inches(5.2), Inches(10), Inches(0.8),
    "\u0413\u0414\u041f\u0411\u0417\u041d \u2014 \u041e\u043f\u0435\u0440\u0430\u0442\u0438\u0432\u0435\u043d \u0446\u0435\u043d\u0442\u044a\u0440 \u0437\u0430 \u0443\u043f\u0440\u0430\u0432\u043b\u0435\u043d\u0438\u0435 \u043d\u0430 \u043f\u0440\u043e\u0438\u0437\u0448\u0435\u0441\u0442\u0432\u0438\u044f",
    size=16, color=GRAY)

prs.save("gdpbzn-presentation.pptx")
print(f"Presentation saved: gdpbzn-presentation.pptx")
print(f"Slides: {len(prs.slides)}")
